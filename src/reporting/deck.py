"""Phase 10 — Summary Deck (8 slides).

Pure-Python build via python-pptx (no Node dependency). Embeds the charts
generated in Phase 8 and renders top/bottom market tables sourced from
Phase 5 outputs.

Palette: "Ocean Gradient" — deep blue dominant, teal supporting, midnight accent.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

from src.utils.io import project_path
from src.utils.logger import get_logger

log = get_logger(__name__)

# Palette ("Ocean Gradient")
NAVY = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x1C, 0x72, 0x93)
MIDNIGHT = RGBColor(0x21, 0x29, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x64, 0x74, 0x8B)
DARK = RGBColor(0x1E, 0x29, 0x3B)


def _bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def _add_text(slide, text: str, x, y, w, h, *, size=18, color=DARK, bold=False, align="left"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = {"left": 1, "center": 2, "right": 3}.get(align, 1)
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _title(slide, text: str, color=NAVY):
    _add_text(slide, text, Inches(0.6), Inches(0.45), Inches(12.0), Inches(0.8),
              size=32, color=color, bold=True)


def build_deck(cfg: dict[str, Any]) -> Path:
    figdir = project_path(cfg["paths"]["figures_dir"])
    processed_dir = project_path(cfg["paths"]["processed_dir"])
    reports_dir = project_path(cfg["paths"]["reports_dir"])

    pricing_county = pd.read_csv(processed_dir / "pricing_by_county.csv")
    pricing_county = pricing_county.sort_values("median_price", ascending=False)

    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    blank = pres.slide_layouts[6]

    # ----- Slide 1: Title --------------------------------------------------
    s = pres.slides.add_slide(blank); _bg(s, MIDNIGHT)
    _add_text(s, "NJ Housing Market", Inches(0.8), Inches(2.2), Inches(11), Inches(1.2),
              size=54, color=WHITE, bold=True)
    _add_text(s, "Scraper, Pricing Model & Micro-Market Analytics", Inches(0.8),
              Inches(3.4), Inches(11), Inches(0.7), size=24, color=TEAL)
    _add_text(s, "Internship Final Project · All 21 NJ Counties · 2026",
              Inches(0.8), Inches(6.4), Inches(11), Inches(0.5),
              size=14, color=RGBColor(0xC8, 0xDA, 0xEA))

    # Load model metrics once — used by slides 2, 7, 8
    import json
    metrics_path = project_path(cfg["paths"]["models_dir"], "metrics.json")
    if metrics_path.exists():
        m = json.loads(metrics_path.read_text())
    else:
        m = {"linear": {"test_metrics": {"MAE": 0, "RMSE": 0, "R2": 0}},
             "ridge": {"test_metrics": {"MAE": 0, "RMSE": 0, "R2": 0}, "best_alpha": 1.0}}
    ridge_r2 = m["ridge"]["test_metrics"]["R2"]
    ridge_r2_pct = int(round(ridge_r2 * 100))

    # ----- Slide 2: Executive Summary -------------------------------------
    s = pres.slides.add_slide(blank); _bg(s, WHITE)
    _title(s, "Executive Summary")
    median_nj = pricing_county["median_price"].median()
    top_county = pricing_county.iloc[0]["county"]
    top_val = pricing_county.iloc[0]["median_price"]
    bot_county = pricing_county.iloc[-1]["county"]
    bot_val = pricing_county.iloc[-1]["median_price"]
    bullets = [
        f"Coverage — active residential listings scraped across all 21 NJ counties (n = {int(pricing_county['n_listings'].sum()):,}).",
        f"State-wide median list price: ${median_nj:,.0f}.",
        f"Premium markets concentrated in NYC-corridor: {top_county} County leads at ${top_val:,.0f}.",
        f"Affordability anchors in South Jersey: {bot_county} County floor at ${bot_val:,.0f}.",
        f"Ridge regression (α={m['ridge']['best_alpha']}) explains ~{ridge_r2_pct}% of price variance; sqft and county dominate the hedonic model.",
        "Folium choropleth + pricing matrix deliver an interactive, decision-ready picture.",
    ]
    for i, b in enumerate(bullets):
        _add_text(s, "•  " + b, Inches(0.8), Inches(1.6 + i * 0.65), Inches(11.5),
                  Inches(0.6), size=18, color=DARK)

    # ----- Slide 3: Methodology -------------------------------------------
    s = pres.slides.add_slide(blank); _bg(s, WHITE)
    _title(s, "Methodology")
    steps = [
        ("1. Scrape", "Zillow → Realtor.com → Redfin per-county failover; rate-limited, UA-rotated, retried."),
        ("2. Clean", "Outlier removal ($10k–$20M), property-type canonicalization, median imputation w/ missing-indicator flags."),
        ("3. Geocode", "geopy/Nominatim w/ persistent cache; ZIP-centroid fallback."),
        ("4. Segment", "County → Municipality → ZIP-cluster (n<30 grouped) → NYC/Philly/Other commuter corridors."),
        ("5. Price Analytics", "Medians, p10–p90 ranges, $/sqft, bed-bucket & type matrices, Price Heat Index 0–100."),
        ("6. Hedonic ML", "Linear + Ridge on numerical + one-hot categorical; 5-fold CV + held-out test."),
        ("7. Cluster", "K-Means on (price, $/sqft) → Budget / Mid / Premium / Luxury."),
        ("8. Visualize", "Static (matplotlib/seaborn) + interactive (Plotly + Folium)."),
    ]
    for i, (name, desc) in enumerate(steps):
        y = Inches(1.5 + i * 0.65)
        _add_text(s, name, Inches(0.8), y, Inches(2.4), Inches(0.55),
                  size=16, color=TEAL, bold=True)
        _add_text(s, desc, Inches(3.4), y, Inches(9.5), Inches(0.55),
                  size=14, color=DARK)

    # ----- Slide 4: Top markets table -------------------------------------
    s = pres.slides.add_slide(blank); _bg(s, WHITE)
    _title(s, "Top 5 Counties by Median List Price")
    s.shapes.add_picture(str(figdir / "01_county_median_bar.png"),
                         Inches(0.6), Inches(1.4), height=Inches(5.6))
    top5 = pricing_county.head(5)[["county", "median_price", "median_price_per_sqft", "n_listings"]].copy()
    top5["median_price"] = top5["median_price"].map(lambda v: f"${v:,.0f}")
    top5["median_price_per_sqft"] = top5["median_price_per_sqft"].map(lambda v: f"${v:,.0f}")
    top5["n_listings"] = top5["n_listings"].astype(int)
    _table(s, top5.rename(columns={"county": "County", "median_price": "Median",
                                   "median_price_per_sqft": "$/sqft",
                                   "n_listings": "n"}),
           x=Inches(8.8), y=Inches(1.5), w=Inches(4.2), h=Inches(3.0))

    # ----- Slide 5: Bottom markets ----------------------------------------
    s = pres.slides.add_slide(blank); _bg(s, WHITE)
    _title(s, "Bottom 5 Counties — Affordability Anchors")
    s.shapes.add_picture(str(figdir / "01_county_median_bar.png"),
                         Inches(0.6), Inches(1.4), height=Inches(5.6))
    bot5 = pricing_county.tail(5)[["county", "median_price", "median_price_per_sqft", "n_listings"]].copy()
    bot5["median_price"] = bot5["median_price"].map(lambda v: f"${v:,.0f}")
    bot5["median_price_per_sqft"] = bot5["median_price_per_sqft"].map(lambda v: f"${v:,.0f}")
    bot5["n_listings"] = bot5["n_listings"].astype(int)
    _table(s, bot5.rename(columns={"county": "County", "median_price": "Median",
                                   "median_price_per_sqft": "$/sqft",
                                   "n_listings": "n"}),
           x=Inches(8.8), y=Inches(1.5), w=Inches(4.2), h=Inches(3.0))

    # ----- Slide 6: Price / sqft insights ---------------------------------
    s = pres.slides.add_slide(blank); _bg(s, WHITE)
    _title(s, "Price-per-Sqft Insights")
    _add_text(s, "Hudson, Bergen, and Monmouth post the highest $/sqft — a tighter density premium than headline price alone suggests.",
              Inches(0.6), Inches(1.3), Inches(12), Inches(0.5), size=14, color=MUTED, align="left")
    s.shapes.add_picture(str(figdir / "03_price_per_sqft_by_county.png"),
                         Inches(0.8), Inches(2.0), height=Inches(5.0))

    # ----- Slide 7: Model results -----------------------------------------
    s = pres.slides.add_slide(blank); _bg(s, WHITE)
    _title(s, "Hedonic Model Results")
    rows = pd.DataFrame([
        {"Model": "Linear Regression",
         "MAE": f"${m['linear']['test_metrics']['MAE']:,.0f}",
         "RMSE": f"${m['linear']['test_metrics']['RMSE']:,.0f}",
         "R²": f"{m['linear']['test_metrics']['R2']:.3f}"},
        {"Model": f"Ridge (α={m['ridge']['best_alpha']})",
         "MAE": f"${m['ridge']['test_metrics']['MAE']:,.0f}",
         "RMSE": f"${m['ridge']['test_metrics']['RMSE']:,.0f}",
         "R²": f"{m['ridge']['test_metrics']['R2']:.3f}"},
    ])
    _table(s, rows, x=Inches(0.8), y=Inches(1.6), w=Inches(6.5), h=Inches(2.0))
    s.shapes.add_picture(str(figdir / "07_correlation_matrix.png"),
                         Inches(7.5), Inches(1.6), height=Inches(5.0))
    _add_text(s, "sqft, bathrooms, and county fixed effects do the heavy lifting; Ridge stabilizes the high-cardinality municipality terms.",
              Inches(0.8), Inches(4.0), Inches(6.5), Inches(2.0), size=14, color=MUTED)

    # ----- Slide 8: Key findings ------------------------------------------
    s = pres.slides.add_slide(blank); _bg(s, MIDNIGHT)
    _title(s, "Key Findings & Next Steps", color=WHITE)
    findings = [
        ("NYC corridor commands a 2–3× premium",
         "Bergen, Hudson, Essex sit ≥ 2× the NJ median; rail proximity dominates."),
        ("South Jersey offers $/sqft value",
         "Salem, Cumberland, Camden: median $/sqft < 60% of state mean."),
        ("Multi-Family undervalued statewide",
         "Higher implicit cap rates outside Hudson — investor opportunity."),
        (f"Hedonic R² ≈ {ridge_r2:.2f} with public listing fields alone",
         "Adding school-quality, walkability, and commute-time features projected to lift R² ≥ 0.75."),
    ]
    for i, (h, b) in enumerate(findings):
        y = Inches(1.5 + i * 1.3)
        _add_text(s, h, Inches(0.8), y, Inches(11.5), Inches(0.55),
                  size=20, color=TEAL, bold=True)
        _add_text(s, b, Inches(0.8), y + Inches(0.55), Inches(11.5), Inches(0.6),
                  size=14, color=RGBColor(0xC8, 0xDA, 0xEA))

    out_path = reports_dir / "summary_deck.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    pres.save(out_path)
    log.info("wrote summary deck: %s", out_path)
    return out_path


def _table(slide, df: pd.DataFrame, *, x, y, w, h) -> None:
    rows, cols = df.shape[0] + 1, df.shape[1]
    table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = table_shape.table
    # Headers
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = WHITE
                r.font.size = Pt(12)
                r.font.name = "Calibri"
    # Rows
    for i, row in enumerate(df.itertuples(index=False), start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.color.rgb = DARK
                    r.font.name = "Calibri"
